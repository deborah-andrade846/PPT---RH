# -*- coding: utf-8 -*-
"""
Gera a apresentação executiva:
"Fluxo do Processo de Recrutamento e Seleção"

Formato 16:9, design corporativo (azul / amarelo / cinza),
com fluxograma (caixas de atividade, losangos de decisão,
conectores), painéis informativos e linha do tempo executiva.

Dependência: python-pptx  ->  pip install python-pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from pptx.enum.text import MSO_AUTO_SIZE
import copy

# --------------------------------------------------------------------------
# PALETA CORPORATIVA  (azul / amarelo / cinza)
# --------------------------------------------------------------------------
NAVY      = "16305C"   # azul escuro / títulos
BLUE      = "1F4E9B"   # azul principal
BLUE_MED  = "2E6CC4"   # azul médio
BLUE_SOFT = "DCE7F7"   # azul claro (preenchimento)
BLUE_PALE = "EEF3FB"   # azul quase branco
YELLOW    = "F5B301"   # amarelo destaque
YELLOW_DK = "D99400"   # amarelo escuro
YELLOW_SF = "FCE9BE"   # amarelo suave
GRAY_DK   = "3A414C"   # texto principal
GRAY      = "6B7480"   # texto secundário
GRAY_LT   = "D7DCE4"   # linhas / bordas
GRAY_BG   = "F4F6F9"   # fundo suave
WHITE     = "FFFFFF"
STOP      = "C0504D"   # encerramento (vermelho sóbrio)
STOP_SF   = "F4DAD8"

FONT = "Calibri"

EMU_IN = 914400

# --------------------------------------------------------------------------
# HELPERS
# --------------------------------------------------------------------------
def _rgb(hexstr):
    return RGBColor.from_string(hexstr)


def solid_fill(shape, hexstr):
    shape.fill.solid()
    shape.fill.fore_color.rgb = _rgb(hexstr)


def no_line(shape):
    shape.line.fill.background()


def set_line(shape, hexstr, width_pt=1.0):
    shape.line.color.rgb = _rgb(hexstr)
    shape.line.width = Pt(width_pt)


def add_shadow(shape, blur=0.07, dist=0.045, direction=5400000, alpha=72):
    """Sombra externa suave (soft drop shadow)."""
    spPr = shape._element.spPr
    # remove efeito herdado
    for tag in ("a:effectLst",):
        existing = spPr.find(qn(tag))
        if existing is not None:
            spPr.remove(existing)
    effectLst = spPr.makeelement(qn("a:effectLst"), {})
    shdw = effectLst.makeelement(qn("a:outerShdw"), {
        "blurRad": str(int(blur * EMU_IN)),
        "dist": str(int(dist * EMU_IN)),
        "dir": str(direction),
        "rotWithShape": "0",
    })
    clr = shdw.makeelement(qn("a:srgbClr"), {"val": "1B2A44"})
    a = clr.makeelement(qn("a:alpha"), {"val": str(int((100 - alpha) * 1000 + alpha * 0))})
    # alpha here: opacity percentage *1000. Lower = more transparent.
    a.set("val", str(int(alpha * 1000 // 3)))
    clr.append(a)
    shdw.append(clr)
    effectLst.append(shdw)
    spPr.append(effectLst)


def _set_text(shape, blocks, anchor=MSO_ANCHOR.MIDDLE, wrap=True,
              l=0.06, r=0.06, t=0.02, b=0.02):
    """
    blocks: lista de dicts com chaves:
      text, size, bold, color, align, space_after, space_before, italic
    Cada dict é um parágrafo. Runs simples.
    """
    tf = shape.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    try:
        tf.auto_size = MSO_AUTO_SIZE.NONE
    except Exception:
        pass
    tf.margin_left = Inches(l)
    tf.margin_right = Inches(r)
    tf.margin_top = Inches(t)
    tf.margin_bottom = Inches(b)

    for i, blk in enumerate(blocks):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = blk.get("align", PP_ALIGN.CENTER)
        if "space_after" in blk:
            p.space_after = Pt(blk["space_after"])
        else:
            p.space_after = Pt(1)
        if "space_before" in blk:
            p.space_before = Pt(blk["space_before"])
        else:
            p.space_before = Pt(0)
        if blk.get("line_spacing"):
            p.line_spacing = blk["line_spacing"]
        run = p.add_run()
        run.text = blk["text"]
        f = run.font
        f.name = FONT
        f.size = Pt(blk.get("size", 11))
        f.bold = blk.get("bold", False)
        f.italic = blk.get("italic", False)
        f.color.rgb = _rgb(blk.get("color", GRAY_DK))


def add_rect(slide, x, y, w, h, fill=WHITE, line=None, line_w=1.0,
             rounded=True, radius=0.09, shadow=True):
    shp_type = MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE
    shp = slide.shapes.add_shape(shp_type, Inches(x), Inches(y),
                                 Inches(w), Inches(h))
    if rounded:
        try:
            shp.adjustments[0] = radius
        except Exception:
            pass
    solid_fill(shp, fill)
    if line:
        set_line(shp, line, line_w)
    else:
        no_line(shp)
    if shadow:
        add_shadow(shp)
    shp.shadow.inherit = False if not shadow else shp.shadow.inherit
    return shp


def add_diamond(slide, cx, cy, w, h, fill=YELLOW, line=YELLOW_DK,
                text="", shadow=True):
    x = cx - w / 2
    y = cy - h / 2
    shp = slide.shapes.add_shape(MSO_SHAPE.DIAMOND, Inches(x), Inches(y),
                                 Inches(w), Inches(h))
    solid_fill(shp, fill)
    if line:
        set_line(shp, line, 1.0)
    else:
        no_line(shp)
    if shadow:
        add_shadow(shp)
    if text:
        _set_text(shp, [{"text": text, "size": 9, "bold": True,
                         "color": NAVY, "align": PP_ALIGN.CENTER}])
    return shp


def add_arrow(slide, x1, y1, x2, y2, color=BLUE_MED, width=1.6,
              connector=MSO_CONNECTOR.STRAIGHT, head=True):
    conn = slide.shapes.add_connector(connector, Inches(x1), Inches(y1),
                                      Inches(x2), Inches(y2))
    conn.line.color.rgb = _rgb(color)
    conn.line.width = Pt(width)
    ln = conn.line._get_or_add_ln()
    if head:
        tail = ln.makeelement(qn("a:tailEnd"),
                              {"type": "triangle", "w": "med", "len": "med"})
        ln.append(tail)
    return conn


def add_text(slide, x, y, w, h, blocks, anchor=MSO_ANCHOR.MIDDLE,
             wrap=True, l=0.05, r=0.05, t=0.02, b=0.02):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    _set_text(tb, blocks, anchor=anchor, wrap=wrap, l=l, r=r, t=t, b=b)
    return tb


def add_chip(slide, cx, cy, d, glyph, chip_color=YELLOW, glyph_color=NAVY,
             circle=True, gsize=11):
    """Pequeno 'ícone' circular com um glifo/numero."""
    x = cx - d / 2
    y = cy - d / 2
    shp_type = MSO_SHAPE.OVAL if circle else MSO_SHAPE.ROUNDED_RECTANGLE
    shp = slide.shapes.add_shape(shp_type, Inches(x), Inches(y),
                                 Inches(d), Inches(d))
    solid_fill(shp, chip_color)
    no_line(shp)
    add_shadow(shp, blur=0.04, dist=0.03, alpha=55)
    _set_text(shp, [{"text": glyph, "size": gsize, "bold": True,
                     "color": glyph_color, "align": PP_ALIGN.CENTER}],
              l=0, r=0, t=0, b=0)
    return shp


# ==========================================================================
# APRESENTAÇÃO
# ==========================================================================
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
W, H = 13.333, 7.5


def bg(slide, color=WHITE):
    r = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0,
                               prs.slide_width, prs.slide_height)
    solid_fill(r, color)
    no_line(r)
    r.shadow.inherit = False
    # envia para trás
    sp = r._element
    sp.getparent().remove(sp)
    slide.shapes._spTree.insert(2, sp)
    return r


# --------------------------------------------------------------------------
# SLIDE 1 — TÍTULO
# --------------------------------------------------------------------------
s1 = prs.slides.add_slide(BLANK)
bg(s1, WHITE)

# Painel lateral esquerdo (azul)
panel = add_rect(s1, 0, 0, 4.9, 7.5, fill=NAVY, rounded=False, shadow=False)
# faixa amarela vertical de acento
add_rect(s1, 4.9, 0, 0.12, 7.5, fill=YELLOW, rounded=False, shadow=False)

# Elementos gráficos decorativos (círculos)
for (cx, cy, d, col, al) in [(4.9, 0.0, 3.0, BLUE, None),
                             (0.2, 6.9, 2.4, BLUE_MED, None)]:
    c = s1.shapes.add_shape(MSO_SHAPE.OVAL, Inches(cx - d/2), Inches(cy - d/2),
                            Inches(d), Inches(d))
    solid_fill(c, col)
    no_line(c)
    c.shadow.inherit = False

# --- Motivo gráfico de RH: rede de pessoas (nós ligados) ---
def person(slide, cx, cy, scale=1.0, color=WHITE):
    head_d = 0.30 * scale
    body_w = 0.54 * scale
    body_h = 0.34 * scale
    hd = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                Inches(cx - head_d/2), Inches(cy - head_d/2),
                                Inches(head_d), Inches(head_d))
    solid_fill(hd, color); no_line(hd); hd.shadow.inherit = False
    bd = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                Inches(cx - body_w/2), Inches(cy + head_d/2 - 0.02*scale),
                                Inches(body_w), Inches(body_h))
    try:
        bd.adjustments[0] = 0.5
    except Exception:
        pass
    solid_fill(bd, color); no_line(bd); bd.shadow.inherit = False

nodes = [(2.45, 2.55, YELLOW, 1.25),
         (1.35, 3.65, WHITE, 1.0),
         (3.55, 3.65, WHITE, 1.0),
         (2.45, 4.75, BLUE_MED, 1.05)]
# linhas de conexão entre pessoas
for i in range(len(nodes)):
    for j in range(i + 1, len(nodes)):
        add_arrow(s1, nodes[i][0], nodes[i][1] + 0.15,
                  nodes[j][0], nodes[j][1] + 0.15,
                  color="3E5C8A", width=1.2, head=False)
for (cx, cy, col, sc) in nodes:
    person(s1, cx, cy, scale=sc, color=col)

# Rótulo superior
add_text(s1, 0.55, 0.55, 4.0, 0.4,
         [{"text": "RECURSOS HUMANOS", "size": 12, "bold": True,
           "color": YELLOW, "align": PP_ALIGN.LEFT}],
         anchor=MSO_ANCHOR.TOP)

# Título e subtítulo (lado direito)
add_text(s1, 5.45, 2.25, 7.3, 1.9,
         [{"text": "Fluxo do Processo de", "size": 34, "bold": True,
           "color": NAVY, "align": PP_ALIGN.LEFT, "line_spacing": 1.02},
          {"text": "Recrutamento e Seleção", "size": 34, "bold": True,
           "color": BLUE, "align": PP_ALIGN.LEFT, "line_spacing": 1.02}],
         anchor=MSO_ANCHOR.BOTTOM)

# linha amarela sob o título
add_rect(s1, 5.5, 4.25, 1.7, 0.06, fill=YELLOW, rounded=False, shadow=False)

add_text(s1, 5.45, 4.45, 7.2, 0.9,
         [{"text": "Etapas, aprovações, sistemas utilizados e pontos de decisão",
           "size": 15, "bold": False, "color": GRAY, "align": PP_ALIGN.LEFT}],
         anchor=MSO_ANCHOR.TOP)

# rodapé pequeno
add_text(s1, 5.45, 6.75, 7.2, 0.4,
         [{"text": "Apresentação executiva  ·  Processo de R&S",
           "size": 10, "bold": True, "color": GRAY, "align": PP_ALIGN.LEFT}],
         anchor=MSO_ANCHOR.TOP)


# --------------------------------------------------------------------------
# SLIDE 2 — FLUXO COMPLETO DO PROCESSO
# --------------------------------------------------------------------------
s2 = prs.slides.add_slide(BLANK)
bg(s2, GRAY_BG)

# Cabeçalho
add_rect(s2, 0, 0, 13.333, 0.86, fill=NAVY, rounded=False, shadow=False)
add_rect(s2, 0, 0.86, 13.333, 0.06, fill=YELLOW, rounded=False, shadow=False)
add_text(s2, 0.4, 0.0, 9.0, 0.86,
         [{"text": "Fluxo Completo do Processo de Recrutamento e Seleção",
           "size": 20, "bold": True, "color": WHITE, "align": PP_ALIGN.LEFT}],
         anchor=MSO_ANCHOR.MIDDLE)
add_text(s2, 9.4, 0.0, 3.6, 0.86,
         [{"text": "Fluxograma executivo", "size": 11, "bold": True,
           "color": YELLOW, "align": PP_ALIGN.RIGHT}],
         anchor=MSO_ANCHOR.MIDDLE)


def stage_box(slide, x, y, w, h, num, title, subs, fill=WHITE,
              title_color=NAVY, bar=BLUE, chip=BLUE):
    """Caixa de atividade com barra lateral, número e sub-bullets."""
    box = add_rect(slide, x, y, w, h, fill=fill, rounded=True, radius=0.10,
                   shadow=True)
    # barra lateral colorida
    add_rect(slide, x, y, 0.10, h, fill=bar, rounded=False, shadow=False)
    # chip numerado
    add_chip(slide, x + 0.42, y + 0.35, 0.44, str(num),
             chip_color=chip, glyph_color=WHITE, gsize=13)
    blocks = [{"text": title, "size": 11.5, "bold": True,
               "color": title_color, "align": PP_ALIGN.LEFT,
               "space_after": 2}]
    for sline in subs:
        blocks.append({"text": sline, "size": 8, "bold": False,
                       "color": GRAY, "align": PP_ALIGN.LEFT,
                       "space_after": 0, "line_spacing": 1.0})
    tb = slide.shapes.add_textbox(Inches(x + 0.72), Inches(y + 0.03),
                                  Inches(w - 0.82), Inches(h - 0.06))
    _set_text(tb, blocks, anchor=MSO_ANCHOR.MIDDLE, l=0.02, r=0.04,
              t=0.02, b=0.02)
    return box


def stop_node(slide, x, y, w, text):
    b = add_rect(slide, x, y, w, 0.42, fill=STOP_SF, line=STOP, line_w=1.0,
                 rounded=True, radius=0.28, shadow=True)
    _set_text(b, [{"text": text, "size": 8, "bold": True,
                   "color": STOP, "align": PP_ALIGN.CENTER}])
    return b


def yn_label(slide, x, y, text, color):
    add_text(slide, x, y, 0.55, 0.24,
             [{"text": text, "size": 8.5, "bold": True, "color": color,
               "align": PP_ALIGN.CENTER}], anchor=MSO_ANCHOR.MIDDLE,
             l=0.01, r=0.01, t=0, b=0)


# ---- LANE 1 (fluxo principal, coluna esquerda) ----
LX = 0.45          # x da coluna 1
LW = 3.05          # largura
cx1 = LX + LW / 2

# 1 Requisição de Vaga
stage_box(s2, LX, 1.10, LW, 0.90, 1, "Requisição de Vaga",
          ["Origem: Portal de Serviços (HUB RH)",
           "Gestor solicita abertura da vaga",
           "Cargo · Centro de custo · Escopo · Contratação"],
          bar=BLUE, chip=BLUE)
# 2 Publicação da Vaga
stage_box(s2, LX, 2.18, LW, 0.86, 2, "Publicação da Vaga",
          ["Publicação na Pandapé",
           "Quintas-feiras · permanência de 7 dias",
           "Portal de Vagas · Aura Comunica (sexta)"],
          bar=BLUE, chip=BLUE)
# 3 Triagem de Currículos
stage_box(s2, LX, 3.22, LW, 0.80, 3, "Triagem de Currículos",
          ["Avaliação técnica dos currículos",
           "Compatibilidade com os requisitos",
           "Seleção de candidatos aderentes"],
          bar=BLUE, chip=BLUE)
# 4 Entrevista RH
stage_box(s2, LX, 4.20, LW, 0.86, 4, "Entrevista RH",
          ["Avaliação comportamental · histórico",
           "Motivação · alinhamento cultural",
           "Pretensão salarial"],
          bar=BLUE, chip=BLUE)

# conectores lane 1
for (ya, yb) in [(2.00, 2.18), (3.04, 3.22), (4.02, 4.20)]:
    add_arrow(s2, cx1, ya, cx1, yb, color=BLUE_MED, width=1.8)

# Diamante D1 - Aprovado? (RH)
d1cy = 5.55
add_arrow(s2, cx1, 5.06, cx1, d1cy - 0.42, color=BLUE_MED, width=1.8)
add_diamond(s2, cx1, d1cy, 1.55, 0.85, fill=YELLOW, line=YELLOW_DK,
            text="Aprovado?")
# Não -> encerramento (à esquerda/abaixo)
yn_label(s2, LX - 0.02, d1cy - 0.32, "Não", STOP)
add_arrow(s2, cx1 - 0.78, d1cy, LX + 0.25, d1cy, color=STOP, width=1.5)
stop_node(s2, LX - 0.35, d1cy + 0.30, 1.55, "Reprovado · processo encerrado")
add_arrow(s2, LX + 0.25, d1cy + 0.02, LX + 0.42, d1cy + 0.30,
          color=STOP, width=1.4, head=False)
# Sim -> inicia a Triagem Patrimonial (entra no topo do branch)
yn_label(s2, cx1 + 0.42, d1cy - 0.02, "Sim", "2E7D32")

# ---- BRANCH PARALELO: Triagem Patrimonial (centro) ----
PX = 3.95
PW = 2.05
pcx = PX + PW / 2

# cabeçalho do branch
hdr = add_rect(s2, PX, 1.10, PW, 0.44, fill=YELLOW, rounded=True,
               radius=0.22, shadow=True)
_set_text(hdr, [{"text": "Triagem Patrimonial", "size": 10.5, "bold": True,
                 "color": NAVY, "align": PP_ALIGN.CENTER}])
add_text(s2, PX, 1.56, PW, 0.24,
         [{"text": "fluxo paralelo", "size": 8, "italic": True,
           "color": GRAY, "align": PP_ALIGN.CENTER}], anchor=MSO_ANCHOR.TOP)

def mini(slide, y, text):
    b = add_rect(slide, PX, y, PW, 0.46, fill=WHITE, line=GRAY_LT, line_w=0.75,
                 rounded=True, radius=0.16, shadow=True)
    _set_text(b, [{"text": text, "size": 8.7, "bold": True,
                   "color": GRAY_DK, "align": PP_ALIGN.CENTER}])
    return b

mini(s2, 1.92, "Envio de Forms")
mini(s2, 2.56, "Preenchimento pelo candidato")
mini(s2, 3.20, "Pesquisa patrimonial")
for (ya, yb) in [(2.38, 2.56), (3.02, 3.20)]:
    add_arrow(s2, pcx, ya, pcx, yb, color=YELLOW_DK, width=1.6)
add_arrow(s2, pcx, 1.80, pcx, 1.92, color=YELLOW_DK, width=1.6)

# diamante patrimonial
pdcy = 4.25
add_arrow(s2, pcx, 3.66, pcx, pdcy - 0.42, color=YELLOW_DK, width=1.6)
add_diamond(s2, pcx, pdcy, 1.55, 0.82, fill=BLUE_SOFT, line=BLUE_MED,
            text="Aprovado?")
# Não -> restrição
yn_label(s2, pcx - 0.5, pdcy + 0.42, "Não", STOP)
add_arrow(s2, pcx, pdcy + 0.41, pcx, pdcy + 0.62, color=STOP, width=1.4)
stop_node(s2, PX, pdcy + 0.62, PW, "Restrição · processo encerrado")

# Ligação: D1 "Sim" -> sobe pelo canal (x=3.72) e entra no topo do branch
CH1 = 3.72
add_arrow(s2, cx1 + 0.78, d1cy, CH1, d1cy, color="2E7D32", width=1.6,
          head=False)
add_arrow(s2, CH1, d1cy, CH1, 1.32, color="2E7D32", width=1.6, head=False)
add_arrow(s2, CH1, 1.32, PX, 1.32, color="2E7D32", width=1.6, head=True)

# Sim (patrimonial) -> Entrevista com Gestor (lane 2)
yn_label(s2, pcx + 0.60, pdcy - 0.30, "Sim", "2E7D32")

# ---- LANE 2 (continuação, coluna direita) ----
QX = 6.35
QW = 3.05
qcx = QX + QW / 2

# 5 Entrevista com Gestor
stage_box(s2, QX, 1.10, QW, 0.86, 5, "Entrevista com Gestor",
          ["Conhecimento técnico · experiência",
           "Competências · fit com a equipe"],
          bar=BLUE, chip=BLUE)
# ligação patrimonial "Sim" -> Entrevista com Gestor (sobe pelo canal x=6.12)
CH2 = 6.12
add_arrow(s2, pcx + 0.78, pdcy, CH2, pdcy, color="2E7D32", width=1.6,
          head=False)
add_arrow(s2, CH2, pdcy, CH2, 1.53, color="2E7D32", width=1.6, head=False)
add_arrow(s2, CH2, 1.53, QX, 1.53, color="2E7D32", width=1.6)

# Diamante D2 - Candidato Finalista?
d2cy = 2.55
add_arrow(s2, qcx, 1.96, qcx, d2cy - 0.41, color=BLUE_MED, width=1.8)
add_diamond(s2, qcx, d2cy, 1.7, 0.82, fill=YELLOW, line=YELLOW_DK,
            text="Candidato\nFinalista?")
# Não -> retorna ao fluxo (nota no rodapé central, via canal x=6.28)
CH3 = 6.28
yn_label(s2, qcx - 1.05, d2cy - 0.02, "Não", STOP)
add_arrow(s2, qcx - 0.85, d2cy, CH3, d2cy, color=STOP, width=1.4, head=False)
add_arrow(s2, CH3, d2cy, CH3, 6.58, color=STOP, width=1.4, head=False)
add_arrow(s2, CH3, 6.58, 6.62, 6.58, color=STOP, width=1.4, head=True)
b = add_rect(s2, 6.62, 6.30, 2.78, 0.56, fill=STOP_SF,
             line=STOP, line_w=1.0, rounded=True, radius=0.16, shadow=True)
_set_text(b, [{"text": "Não finalista  →  retorna à nova triagem "
                       "ou nova publicação da vaga",
               "size": 8, "bold": True, "color": STOP,
               "align": PP_ALIGN.CENTER}])
# Sim
yn_label(s2, qcx + 0.10, d2cy + 0.42, "Sim", "2E7D32")

# 6 Carta Proposta
stage_box(s2, QX, 3.22, QW, 0.80, 6, "Carta Proposta",
          ["Aprovação RH Custos  ·  Aprovação Gestor",
           "Formalização da oferta ao candidato"],
          bar=YELLOW_DK, chip=YELLOW_DK, title_color=NAVY)
add_arrow(s2, qcx, d2cy + 0.41, qcx, 3.22, color=BLUE_MED, width=1.8)

# 7 Formalização & Assinatura
stage_box(s2, QX, 4.18, QW, 0.86, 7, "Formalização & Assinatura",
          ["Envio por e-mail da proposta",
           "Assinatura eletrônica via GOV.BR",
           "Aceite do candidato"],
          bar=BLUE, chip=BLUE)
add_arrow(s2, qcx, 4.02, qcx, 4.18, color=BLUE_MED, width=1.8)

# 8 Seleção Concluída / Onboarding
stage_box(s2, QX, 5.22, QW, 0.90, 8, "Seleção Concluída",
          ["Onboarding do novo colaborador",
           "Integração à equipe e à empresa"],
          bar="2E7D32", chip="2E7D32", title_color="1B5E20", fill="EAF5EC")
add_arrow(s2, qcx, 5.04, qcx, 5.22, color="2E7D32", width=1.8)

# ---- PAINÉIS INFORMATIVOS (direita) ----
def panel(slide, x, y, w, h, title, items, accent, title_bg):
    box = add_rect(slide, x, y, w, h, fill=WHITE, rounded=True, radius=0.06,
                   shadow=True)
    add_rect(slide, x, y, w, 0.40, fill=title_bg, rounded=True, radius=0.10,
             shadow=False)
    # cobre canto inferior do header para dar efeito de barra
    add_rect(slide, x, y + 0.22, w, 0.20, fill=title_bg, rounded=False,
             shadow=False)
    _th = slide.shapes.add_textbox(Inches(x + 0.05), Inches(y),
                                   Inches(w - 0.1), Inches(0.40))
    _set_text(_th, [{"text": title, "size": 10.5, "bold": True,
                     "color": WHITE, "align": PP_ALIGN.CENTER}])
    blocks = []
    for it in items:
        blocks.append({"text": "•  " + it, "size": 8.4, "bold": False,
                       "color": GRAY_DK, "align": PP_ALIGN.LEFT,
                       "space_after": 2.2, "line_spacing": 1.0})
    tb = slide.shapes.add_textbox(Inches(x + 0.14), Inches(y + 0.48),
                                  Inches(w - 0.24), Inches(h - 0.54))
    _set_text(tb, blocks, anchor=MSO_ANCHOR.TOP, l=0.02, r=0.02, t=0.02, b=0.02)
    # marca de acento à esquerda do header
    add_rect(slide, x, y, 0.08, 0.40, fill=accent, rounded=False, shadow=False)

PNX = 9.72
PNW = 3.32
panel(s2, PNX, 1.10, PNW, 2.00,
      "SISTEMAS UTILIZADOS",
      ["Portal de Serviços", "Teams", "Outlook", "Pandapé",
       "Microsoft Forms", "GOV.BR", "Aura Comunica"],
      accent=YELLOW, title_bg=BLUE)
panel(s2, PNX, 3.24, PNW, 1.86,
      "PONTOS DE DECISÃO",
      ["Aprovação da Requisição", "Aprovação RH Custos",
       "Aprovação Entrevista RH", "Aprovação Triagem Patrimonial",
       "Aprovação Gestor", "Aceite da Proposta"],
      accent=YELLOW, title_bg=BLUE)
panel(s2, PNX, 5.24, PNW, 1.72,
      "POSSÍVEIS ENCERRAMENTOS",
      ["Reprovação na Entrevista RH", "Restrição na Triagem Patrimonial",
       "Reprovação pelo Gestor", "Recusa da Proposta",
       "Reabertura da vaga"],
      accent=YELLOW, title_bg=STOP)

# Legenda inferior
leg_y = 7.06
add_text(s2, 0.45, leg_y, 9.0, 0.34,
         [{"text": "Legenda:   ", "size": 8.5, "bold": True, "color": GRAY_DK,
           "align": PP_ALIGN.LEFT}], anchor=MSO_ANCHOR.MIDDLE)
# amostras
def legend_item(x, shp_color, label, diamond=False):
    if diamond:
        add_diamond(s2, x, leg_y + 0.17, 0.26, 0.22, fill=shp_color,
                    line=None, shadow=False)
    else:
        add_rect(s2, x - 0.13, leg_y + 0.06, 0.26, 0.22, fill=shp_color,
                 rounded=True, radius=0.2, shadow=False)
    add_text(s2, x + 0.18, leg_y, 1.4, 0.34,
             [{"text": label, "size": 8, "bold": False, "color": GRAY,
               "align": PP_ALIGN.LEFT}], anchor=MSO_ANCHOR.MIDDLE)

legend_item(1.60, BLUE, "Atividade")
legend_item(3.55, YELLOW, "Decisão", diamond=True)
legend_item(5.50, STOP_SF, "Encerramento")
legend_item(7.55, "EAF5EC", "Conclusão")


# --------------------------------------------------------------------------
# SLIDE 3 — VISÃO EXECUTIVA (linha do tempo)
# --------------------------------------------------------------------------
s3 = prs.slides.add_slide(BLANK)
bg(s3, WHITE)

add_rect(s3, 0, 0, 13.333, 0.86, fill=NAVY, rounded=False, shadow=False)
add_rect(s3, 0, 0.86, 13.333, 0.06, fill=YELLOW, rounded=False, shadow=False)
add_text(s3, 0.4, 0.0, 9.5, 0.86,
         [{"text": "Visão Executiva do Processo", "size": 20, "bold": True,
           "color": WHITE, "align": PP_ALIGN.LEFT}], anchor=MSO_ANCHOR.MIDDLE)
add_text(s3, 9.0, 0.0, 4.0, 0.86,
         [{"text": "Linha do tempo — 8 etapas", "size": 11, "bold": True,
           "color": YELLOW, "align": PP_ALIGN.RIGHT}], anchor=MSO_ANCHOR.MIDDLE)

stages = [
    ("1", "Requisição", "Abertura da vaga pelo gestor"),
    ("2", "Divulgação", "Publicação na Pandapé"),
    ("3", "Triagem", "Análise de currículos"),
    ("4", "Entrevistas", "RH e gestor"),
    ("5", "Aprovação", "Decisões e validações"),
    ("6", "Proposta", "Carta e negociação"),
    ("7", "Admissão", "Assinatura via GOV.BR"),
    ("8", "Onboarding", "Integração do colaborador"),
]

# Linha do tempo horizontal
line_y = 3.85
x_start = 0.95
x_end = 12.4
seg = (x_end - x_start) / (len(stages) - 1)
# trilho
rail = add_rect(s3, x_start, line_y - 0.03, x_end - x_start, 0.06,
                fill=GRAY_LT, rounded=False, shadow=False)

alt_colors = [BLUE, YELLOW_DK]
for i, (num, title, desc) in enumerate(stages):
    cx = x_start + seg * i
    up = (i % 2 == 0)
    node_col = BLUE if up else YELLOW_DK
    # conector vertical
    card_h = 1.32
    if up:
        card_y = line_y - 0.55 - card_h
        add_arrow(s3, cx, line_y - 0.35, cx, card_y + card_h, color=node_col,
                  width=1.4, head=False)
    else:
        card_y = line_y + 0.55
        add_arrow(s3, cx, line_y + 0.35, cx, card_y, color=node_col,
                  width=1.4, head=False)
    # cartão
    cw = 1.42
    card = add_rect(s3, cx - cw/2, card_y, cw, card_h, fill=WHITE,
                    line=GRAY_LT, line_w=0.75, rounded=True, radius=0.10,
                    shadow=True)
    add_rect(s3, cx - cw/2, card_y, cw, 0.09, fill=node_col, rounded=False,
             shadow=False)
    tb = s3.shapes.add_textbox(Inches(cx - cw/2 + 0.06), Inches(card_y + 0.12),
                               Inches(cw - 0.12), Inches(card_h - 0.16))
    _set_text(tb, [{"text": title, "size": 12, "bold": True, "color": NAVY,
                    "align": PP_ALIGN.CENTER, "space_after": 3},
                   {"text": desc, "size": 8.3, "bold": False, "color": GRAY,
                    "align": PP_ALIGN.CENTER, "line_spacing": 1.0}],
              anchor=MSO_ANCHOR.TOP)
    # nó circular na linha
    add_chip(s3, cx, line_y, 0.56, num, chip_color=node_col,
             glyph_color=WHITE, gsize=15)

# Faixa de indicadores inferior
ind_y = 6.35
indicators = [
    ("8", "etapas principais", BLUE),
    ("6", "pontos de decisão", YELLOW_DK),
    ("7", "sistemas integrados", BLUE),
    ("5", "possíveis encerramentos", STOP),
]
iw = 2.85
gap = (13.333 - iw * len(indicators)) / (len(indicators) + 1)
for i, (val, lbl, col) in enumerate(indicators):
    x = gap + i * (iw + gap)
    card = add_rect(s3, x, ind_y, iw, 0.82, fill=GRAY_BG, rounded=True,
                    radius=0.12, shadow=True)
    add_rect(s3, x, ind_y, 0.10, 0.82, fill=col, rounded=False, shadow=False)
    add_text(s3, x + 0.18, ind_y, 0.85, 0.82,
             [{"text": val, "size": 26, "bold": True, "color": col,
               "align": PP_ALIGN.CENTER}], anchor=MSO_ANCHOR.MIDDLE)
    add_text(s3, x + 1.02, ind_y, iw - 1.1, 0.82,
             [{"text": lbl, "size": 11, "bold": True, "color": GRAY_DK,
               "align": PP_ALIGN.LEFT}], anchor=MSO_ANCHOR.MIDDLE)

add_text(s3, 0.4, 5.85, 12.5, 0.4,
         [{"text": "Do pedido à integração — um fluxo estruturado, "
                   "auditável e orientado a decisões.",
           "size": 12, "italic": True, "bold": False, "color": GRAY,
           "align": PP_ALIGN.CENTER}], anchor=MSO_ANCHOR.MIDDLE)


# --------------------------------------------------------------------------
OUT = "Fluxo_Recrutamento_Selecao.pptx"
prs.save(OUT)
print("OK ->", OUT)
