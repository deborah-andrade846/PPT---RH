# -*- coding: utf-8 -*-
"""
Gera a apresentação institucional CONSOLIDADA:
"Processo de Recrutamento, Seleção e Onboarding — Consolidação Multiunidades"

Consolida as apresentações "AS IS" das unidades da Aura em um único documento:
  · Aura Apoena (Almas) — Brasil
  · Aura Borborema (BBR) — Brasil
  · Aura Serra Grande — Brasil
  · Aura Minosa (Honduras) — Honduras
  · Aura Guatemala — Guatemala

O documento distingue claramente:
  · PROCESSO PADRÃO da organização (fluxo institucional)
  · PARTICULARIDADES de cada unidade (com rastreabilidade de origem)
  · DIVERGÊNCIAS / INCONSISTÊNCIAS / CONFLITOS entre as apresentações
  · LACUNAS de informação
  · RESUMO das diferenças e OPORTUNIDADES de padronização

Formato 16:9, design corporativo (azul / amarelo / cinza).
Dependência: python-pptx  ->  pip install python-pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from pptx.enum.text import MSO_AUTO_SIZE

# --------------------------------------------------------------------------
# PALETA CORPORATIVA
# --------------------------------------------------------------------------
NAVY      = "16305C"
BLUE      = "1F4E9B"
BLUE_MED  = "2E6CC4"
BLUE_SOFT = "DCE7F7"
BLUE_PALE = "EEF3FB"
YELLOW    = "F5B301"
YELLOW_DK = "D99400"
YELLOW_SF = "FCE9BE"
GRAY_DK   = "3A414C"
GRAY      = "6B7480"
GRAY_LT   = "D7DCE4"
GRAY_BG   = "F4F6F9"
WHITE     = "FFFFFF"
STOP      = "C0504D"
STOP_SF   = "F4DAD8"
GREEN     = "2E7D32"
GREEN_SF  = "EAF5EC"
TEAL      = "0F7D8C"     # particularidade (destaque neutro)
TEAL_SF   = "DDF0F2"
PURPLE    = "6B4C9A"     # lacuna
PURPLE_SF = "ECE6F5"

FONT = "Calibri"
EMU_IN = 914400

# Cores por unidade (rastreabilidade visual)
UNIT_COLORS = {
    "Apoena":      "1F4E9B",
    "Borborema":   "0F7D8C",
    "Serra Grande":"D99400",
    "Honduras":    "8A3A3A",
    "Guatemala":   "6B4C9A",
}


# --------------------------------------------------------------------------
# HELPERS BÁSICOS
# --------------------------------------------------------------------------
def _rgb(hexstr):
    return RGBColor.from_string(hexstr)


def solid_fill(shape, hexstr):
    shape.fill.solid()
    shape.fill.fore_color.rgb = _rgb(hexstr)


def no_line(shape):
    shape.line.fill.background()


def set_line(shape, hexstr, width_pt=1.0):
    shape.line.color.rgb = _rgb(hexstr)
    shape.line.width = Pt(width_pt)


def add_shadow(shape, blur=0.07, dist=0.045, direction=5400000, alpha=72):
    spPr = shape._element.spPr
    existing = spPr.find(qn("a:effectLst"))
    if existing is not None:
        spPr.remove(existing)
    effectLst = spPr.makeelement(qn("a:effectLst"), {})
    shdw = effectLst.makeelement(qn("a:outerShdw"), {
        "blurRad": str(int(blur * EMU_IN)),
        "dist": str(int(dist * EMU_IN)),
        "dir": str(direction),
        "rotWithShape": "0",
    })
    clr = shdw.makeelement(qn("a:srgbClr"), {"val": "1B2A44"})
    a = clr.makeelement(qn("a:alpha"), {"val": str(int(alpha * 1000 // 3))})
    clr.append(a)
    shdw.append(clr)
    effectLst.append(shdw)
    spPr.append(effectLst)


def _set_text(shape, blocks, anchor=MSO_ANCHOR.MIDDLE, wrap=True,
              l=0.06, r=0.06, t=0.02, b=0.02):
    tf = shape.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    try:
        tf.auto_size = MSO_AUTO_SIZE.NONE
    except Exception:
        pass
    tf.margin_left = Inches(l)
    tf.margin_right = Inches(r)
    tf.margin_top = Inches(t)
    tf.margin_bottom = Inches(b)

    for i, blk in enumerate(blocks):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = blk.get("align", PP_ALIGN.CENTER)
        p.space_after = Pt(blk.get("space_after", 1))
        p.space_before = Pt(blk.get("space_before", 0))
        if blk.get("line_spacing"):
            p.line_spacing = blk["line_spacing"]
        for j, run_spec in enumerate(blk.get("runs", [blk])):
            run = p.add_run()
            run.text = run_spec["text"]
            f = run.font
            f.name = FONT
            f.size = Pt(run_spec.get("size", blk.get("size", 11)))
            f.bold = run_spec.get("bold", blk.get("bold", False))
            f.italic = run_spec.get("italic", blk.get("italic", False))
            f.color.rgb = _rgb(run_spec.get("color", blk.get("color", GRAY_DK)))


def add_rect(slide, x, y, w, h, fill=WHITE, line=None, line_w=1.0,
             rounded=True, radius=0.09, shadow=True):
    shp_type = MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE
    shp = slide.shapes.add_shape(shp_type, Inches(x), Inches(y),
                                 Inches(w), Inches(h))
    if rounded:
        try:
            shp.adjustments[0] = radius
        except Exception:
            pass
    solid_fill(shp, fill)
    if line:
        set_line(shp, line, line_w)
    else:
        no_line(shp)
    if shadow:
        add_shadow(shp)
    else:
        shp.shadow.inherit = False
    return shp


def add_diamond(slide, cx, cy, w, h, fill=YELLOW, line=YELLOW_DK,
                text="", shadow=True):
    x = cx - w / 2
    y = cy - h / 2
    shp = slide.shapes.add_shape(MSO_SHAPE.DIAMOND, Inches(x), Inches(y),
                                 Inches(w), Inches(h))
    solid_fill(shp, fill)
    if line:
        set_line(shp, line, 1.0)
    else:
        no_line(shp)
    if shadow:
        add_shadow(shp)
    if text:
        _set_text(shp, [{"text": text, "size": 9, "bold": True,
                         "color": NAVY, "align": PP_ALIGN.CENTER}])
    return shp


def add_arrow(slide, x1, y1, x2, y2, color=BLUE_MED, width=1.6,
              connector=MSO_CONNECTOR.STRAIGHT, head=True):
    conn = slide.shapes.add_connector(connector, Inches(x1), Inches(y1),
                                      Inches(x2), Inches(y2))
    conn.line.color.rgb = _rgb(color)
    conn.line.width = Pt(width)
    ln = conn.line._get_or_add_ln()
    if head:
        tail = ln.makeelement(qn("a:tailEnd"),
                              {"type": "triangle", "w": "med", "len": "med"})
        ln.append(tail)
    return conn


def add_text(slide, x, y, w, h, blocks, anchor=MSO_ANCHOR.MIDDLE,
             wrap=True, l=0.05, r=0.05, t=0.02, b=0.02):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    _set_text(tb, blocks, anchor=anchor, wrap=wrap, l=l, r=r, t=t, b=b)
    return tb


def add_chip(slide, cx, cy, d, glyph, chip_color=YELLOW, glyph_color=NAVY,
             circle=True, gsize=11):
    x = cx - d / 2
    y = cy - d / 2
    shp_type = MSO_SHAPE.OVAL if circle else MSO_SHAPE.ROUNDED_RECTANGLE
    shp = slide.shapes.add_shape(shp_type, Inches(x), Inches(y),
                                 Inches(d), Inches(d))
    solid_fill(shp, chip_color)
    no_line(shp)
    add_shadow(shp, blur=0.04, dist=0.03, alpha=55)
    _set_text(shp, [{"text": glyph, "size": gsize, "bold": True,
                     "color": glyph_color, "align": PP_ALIGN.CENTER}],
              l=0, r=0, t=0, b=0)
    return shp


# ==========================================================================
# APRESENTAÇÃO
# ==========================================================================
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
W, H = 13.333, 7.5


def new_slide(color=WHITE):
    slide = prs.slides.add_slide(BLANK)
    r = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0,
                               prs.slide_width, prs.slide_height)
    solid_fill(r, color)
    no_line(r)
    r.shadow.inherit = False
    sp = r._element
    sp.getparent().remove(sp)
    slide.shapes._spTree.insert(2, sp)
    return slide


def header(slide, title, tag="", subtitle=""):
    add_rect(slide, 0, 0, W, 0.86, fill=NAVY, rounded=False, shadow=False)
    add_rect(slide, 0, 0.86, W, 0.06, fill=YELLOW, rounded=False, shadow=False)
    add_text(slide, 0.4, 0.0, 9.6, 0.86,
             [{"text": title, "size": 20, "bold": True, "color": WHITE,
               "align": PP_ALIGN.LEFT}], anchor=MSO_ANCHOR.MIDDLE)
    if tag:
        add_text(slide, 9.6, 0.0, 3.35, 0.86,
                 [{"text": tag, "size": 11, "bold": True, "color": YELLOW,
                   "align": PP_ALIGN.RIGHT}], anchor=MSO_ANCHOR.MIDDLE)
    if subtitle:
        add_text(slide, 0.42, 0.98, 12.5, 0.34,
                 [{"text": subtitle, "size": 11.5, "italic": True,
                   "color": GRAY, "align": PP_ALIGN.LEFT}],
                 anchor=MSO_ANCHOR.MIDDLE)


def footer(slide, n):
    add_text(slide, 0.4, 7.12, 9.0, 0.3,
             [{"text": "Consolidação Institucional · R&S e Onboarding · Aura",
               "size": 8, "color": GRAY_LT, "align": PP_ALIGN.LEFT}],
             anchor=MSO_ANCHOR.MIDDLE)
    add_text(slide, 12.0, 7.12, 0.9, 0.3,
             [{"text": str(n), "size": 8, "bold": True, "color": GRAY,
               "align": PP_ALIGN.RIGHT}], anchor=MSO_ANCHOR.MIDDLE)


def unit_tag(slide, x, y, unit, w=1.35, h=0.26):
    """Selo de rastreabilidade da unidade de origem."""
    col = UNIT_COLORS.get(unit, GRAY)
    b = add_rect(slide, x, y, w, h, fill=col, rounded=True, radius=0.5,
                 shadow=False)
    _set_text(b, [{"text": unit, "size": 8, "bold": True, "color": WHITE,
                   "align": PP_ALIGN.CENTER}], l=0.03, r=0.03, t=0, b=0)
    return b


def add_table(slide, x, y, w, col_widths, rows, header_fill=NAVY,
              header_color=WHITE, font_size=8.5, header_size=9.0,
              row_h=0.3, zebra=True, first_col_bold=True,
              cell_colors=None, white_text_cells=None):
    """
    rows[0] = cabeçalho. col_widths em polegadas (proporção total = w).
    cell_colors: dict opcional {(r,c): hex} para colorir células específicas.
    white_text_cells: conjunto opcional de (r,c) que devem usar texto branco.
    """
    nrows = len(rows)
    ncols = len(rows[0])
    gtbl = slide.shapes.add_table(nrows, ncols, Inches(x), Inches(y),
                                  Inches(w), Inches(row_h * nrows))
    tbl = gtbl.table
    tbl.first_row = False
    tbl.horz_banding = False
    # remove estilo padrão (borda) — mantemos limpo
    for ci, cw in enumerate(col_widths):
        tbl.columns[ci].width = Inches(cw)
    cell_colors = cell_colors or {}
    white_text_cells = white_text_cells or set()
    for ri, row in enumerate(rows):
        for ci, val in enumerate(row):
            cell = tbl.cell(ri, ci)
            cell.margin_left = Inches(0.06)
            cell.margin_right = Inches(0.05)
            cell.margin_top = Inches(0.02)
            cell.margin_bottom = Inches(0.02)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            if ri == 0:
                fill_hex = header_fill
            elif (ri, ci) in cell_colors:
                fill_hex = cell_colors[(ri, ci)]
            elif zebra and ri % 2 == 0:
                fill_hex = BLUE_PALE
            else:
                fill_hex = WHITE
            cell.fill.solid()
            cell.fill.fore_color.rgb = _rgb(fill_hex)
            tf = cell.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT
            run = p.add_run()
            run.text = str(val)
            f = run.font
            f.name = FONT
            if ri == 0:
                f.size = Pt(header_size)
                f.bold = True
                f.color.rgb = _rgb(header_color)
            else:
                f.size = Pt(font_size)
                f.bold = (ci == 0 and first_col_bold)
                if (ri, ci) in white_text_cells:
                    f.color.rgb = _rgb(WHITE)
                else:
                    f.color.rgb = _rgb(NAVY if (ci == 0 and first_col_bold)
                                       else GRAY_DK)
    return tbl


slide_no = [0]


def finalize(slide):
    slide_no[0] += 1
    footer(slide, slide_no[0])


# ==========================================================================
# SLIDE 1 — CAPA
# ==========================================================================
s = new_slide(WHITE)
add_rect(s, 0, 0, 5.0, 7.5, fill=NAVY, rounded=False, shadow=False)
add_rect(s, 5.0, 0, 0.12, 7.5, fill=YELLOW, rounded=False, shadow=False)
# círculos decorativos
for (cx, cy, d, col) in [(5.0, 0.0, 3.0, BLUE), (0.2, 7.1, 2.4, BLUE_MED)]:
    c = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(cx - d/2), Inches(cy - d/2),
                           Inches(d), Inches(d))
    solid_fill(c, col); no_line(c); c.shadow.inherit = False

add_text(s, 0.55, 0.55, 4.2, 0.4,
         [{"text": "RECURSOS HUMANOS", "size": 12, "bold": True,
           "color": YELLOW, "align": PP_ALIGN.LEFT}], anchor=MSO_ANCHOR.TOP)

# Ícone de unidades conectadas (nós)
nodes = [(2.5, 2.45, YELLOW, 0.34), (1.4, 3.55, WHITE, 0.27),
         (3.6, 3.55, WHITE, 0.27), (1.9, 4.65, BLUE_MED, 0.29),
         (3.1, 4.65, WHITE, 0.27)]
for i in range(len(nodes)):
    for j in range(i + 1, len(nodes)):
        add_arrow(s, nodes[i][0], nodes[i][1], nodes[j][0], nodes[j][1],
                  color="3E5C8A", width=1.1, head=False)
for (cx, cy, col, d) in nodes:
    c = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(cx - d/2), Inches(cy - d/2),
                           Inches(d), Inches(d))
    solid_fill(c, col); no_line(c); c.shadow.inherit = False
add_text(s, 0.55, 5.25, 4.0, 0.4,
         [{"text": "5 unidades · 3 países", "size": 12, "bold": True,
           "color": BLUE_SOFT, "align": PP_ALIGN.LEFT}], anchor=MSO_ANCHOR.TOP)

# Título (lado direito)
add_text(s, 5.5, 1.9, 7.4, 2.2,
         [{"text": "Processo de Recrutamento,", "size": 33, "bold": True,
           "color": NAVY, "align": PP_ALIGN.LEFT, "line_spacing": 1.02},
          {"text": "Seleção e Onboarding", "size": 33, "bold": True,
           "color": BLUE, "align": PP_ALIGN.LEFT, "line_spacing": 1.02}],
         anchor=MSO_ANCHOR.BOTTOM)
add_rect(s, 5.55, 4.2, 1.7, 0.06, fill=YELLOW, rounded=False, shadow=False)
add_text(s, 5.5, 4.4, 7.3, 1.0,
         [{"text": "Consolidação institucional multiunidades", "size": 17,
           "bold": True, "color": GRAY_DK, "align": PP_ALIGN.LEFT},
          {"text": "Processo padrão · particularidades por unidade · "
                   "divergências · lacunas · oportunidades de padronização",
           "size": 12.5, "color": GRAY, "align": PP_ALIGN.LEFT,
           "space_before": 6, "line_spacing": 1.1}],
         anchor=MSO_ANCHOR.TOP)
add_text(s, 5.5, 6.7, 7.3, 0.5,
         [{"text": "Documento executivo  ·  Base: apresentações AS IS das "
                   "unidades  ·  2026", "size": 10, "bold": True,
           "color": GRAY, "align": PP_ALIGN.LEFT}], anchor=MSO_ANCHOR.TOP)
finalize(s)


# ==========================================================================
# SLIDE 2 — CONTEXTO, OBJETIVO E METODOLOGIA
# ==========================================================================
s = new_slide(GRAY_BG)
header(s, "Contexto, Objetivo e Metodologia", tag="Sumário executivo")

def info_card(x, y, w, h, icon, title, items, accent):
    add_rect(s, x, y, w, h, fill=WHITE, rounded=True, radius=0.05, shadow=True)
    add_rect(s, x, y, w, 0.5, fill=accent, rounded=True, radius=0.10,
             shadow=False)
    add_rect(s, x, y + 0.28, w, 0.22, fill=accent, rounded=False, shadow=False)
    add_text(s, x + 0.15, y, w - 0.3, 0.5,
             [{"text": title, "size": 12.5, "bold": True, "color": WHITE,
               "align": PP_ALIGN.LEFT}], anchor=MSO_ANCHOR.MIDDLE)
    blocks = []
    for it in items:
        blocks.append({"text": "•  " + it, "size": 10, "color": GRAY_DK,
                       "align": PP_ALIGN.LEFT, "space_after": 4,
                       "line_spacing": 1.05})
    add_text(s, x + 0.18, y + 0.62, w - 0.36, h - 0.72, blocks,
             anchor=MSO_ANCHOR.TOP)

info_card(0.4, 1.25, 4.05, 2.6, "", "CONTEXTO",
          ["Cada unidade documentou o mesmo processo de R&S e Onboarding "
           "em sua própria apresentação AS IS.",
           "As unidades operam o mesmo processo institucional, mas com "
           "particularidades locais, exceções e regras próprias.",
           "As versões apresentavam redundâncias, divergências e lacunas."],
          BLUE)
info_card(4.65, 1.25, 4.05, 2.6, "", "OBJETIVO",
          ["Consolidar tudo em uma única apresentação fiel ao processo "
           "institucional.",
           "Separar o PADRÃO organizacional das EXCEÇÕES locais.",
           "Preservar todas as informações relevantes, sem perdas.",
           "Evidenciar divergências e lacunas para decisão."],
          TEAL)
info_card(8.9, 1.25, 4.05, 2.6, "", "METODOLOGIA",
          ["Leitura integral das apresentações de todas as unidades.",
           "Identificação do fluxo comum (processo padrão).",
           "Mapeamento de particularidades com rastreabilidade da origem.",
           "Registro de conflitos, inconsistências e lacunas."],
          YELLOW_DK)

# Faixa de escopo (documentos consolidados)
add_rect(s, 0.4, 4.15, 12.53, 2.55, fill=WHITE, rounded=True, radius=0.05,
         shadow=True)
add_rect(s, 0.4, 4.15, 0.12, 2.55, fill=YELLOW, rounded=False, shadow=False)
add_text(s, 0.65, 4.28, 12.0, 0.4,
         [{"text": "Escopo consolidado — 7 apresentações analisadas",
           "size": 13, "bold": True, "color": NAVY, "align": PP_ALIGN.LEFT}],
         anchor=MSO_ANCHOR.MIDDLE)
add_table(s, 0.65, 4.78, 12.0,
          [3.2, 2.1, 1.6, 5.1],
          [["Apresentação (arquivo)", "Unidade", "País", "Conteúdo coberto"],
           ["Fluxo Original / Fluxo_Recrutamento_Selecao", "Aura Apoena (base)",
            "Brasil", "R&S (8 etapas) + Onboarding (10 etapas) — modelo de referência"],
           ["Almas Fluxo_Recrutamento_Selecao", "Aura Apoena (Almas)",
            "Brasil", "R&S + Onboarding, com Conflito de Interesses e Kit Boas-vindas"],
           ["BBR_Fluxo_Recrutamento_Selecao", "Aura Borborema (BBR)",
            "Brasil", "R&S + Onboarding, com Background Check e admissão via Vexia/Pipefy"],
           ["Serra Grande - Crixas_Fluxo...", "Aura Serra Grande",
            "Brasil", "R&S + Onboarding, com Avaliação Comportamental e DP/FPW"],
           ["Honduras - Crixas_Fluxo...", "Aura Minosa",
            "Honduras", "R&S em 3 rotas + Incorporación (10 etapas) — modelo LATAM"],
           ["Reclutamiento y Selección - Guatemala", "Aura Guatemala",
            "Guatemala", "Inducción documentada; R&S em imagens (lacuna de texto)"]],
          row_h=0.265, header_size=9.5, font_size=8.6)
finalize(s)


# ==========================================================================
# SLIDE 3 — COMO LER ESTA APRESENTAÇÃO (GOVERNANÇA / LEGENDA)
# ==========================================================================
s = new_slide(WHITE)
header(s, "Como Ler Esta Apresentação",
       tag="Convenções", subtitle="Quatro camadas de leitura mantêm o padrão "
       "institucional separado das exceções locais e das pendências.")

layers = [
    (BLUE, "PADRÃO", "Processo institucional",
     "Fluxo comum a todas as unidades. Representa o “como deveria ser” "
     "institucional, válido para todas as operações."),
    (TEAL, "PARTICULARIDADE", "Exceção / regra local",
     "Adição, exceção ou regra específica de uma unidade. Sempre acompanhada "
     "do selo da unidade de origem (rastreabilidade)."),
    (STOP, "DIVERGÊNCIA / CONFLITO", "Inconsistência entre versões",
     "Mesma etapa descrita de forma diferente ou conflitante entre unidades "
     "(ou dentro da mesma unidade). Exige decisão."),
    (PURPLE, "LACUNA", "Informação faltante",
     "Ponto do processo não documentado, incompleto ou “em revisão”. "
     "Precisa ser levantado para completar o padrão."),
]
y = 1.5
for col, tag, sub, desc in layers:
    add_rect(s, 0.5, y, 12.3, 1.18, fill=GRAY_BG, rounded=True, radius=0.06,
             shadow=True)
    add_rect(s, 0.5, y, 0.14, 1.18, fill=col, rounded=False, shadow=False)
    chip = add_rect(s, 0.8, y + 0.26, 2.9, 0.66, fill=col, rounded=True,
                    radius=0.16, shadow=False)
    _set_text(chip, [{"text": tag, "size": 13, "bold": True, "color": WHITE,
                      "align": PP_ALIGN.CENTER}])
    add_text(s, 3.95, y + 0.12, 3.0, 0.94,
             [{"text": sub, "size": 12.5, "bold": True, "color": NAVY,
               "align": PP_ALIGN.LEFT}], anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, 7.0, y + 0.10, 5.6, 0.98,
             [{"text": desc, "size": 10.5, "color": GRAY_DK,
               "align": PP_ALIGN.LEFT, "line_spacing": 1.05}],
             anchor=MSO_ANCHOR.MIDDLE)
    y += 1.32
finalize(s)


# ==========================================================================
# SLIDE 4 — PROCESSO PADRÃO: FLUXOGRAMA DE R&S
# ==========================================================================
s = new_slide(GRAY_BG)
header(s, "Processo Padrão — Fluxo de Recrutamento e Seleção",
       tag="PADRÃO institucional",
       subtitle="Fluxo comum às unidades do Brasil (Apoena, Borborema e Serra "
       "Grande). Particularidades e o modelo LATAM são tratados à parte.")

def stage_box(x, y, w, h, num, title, subs, fill=WHITE, title_color=NAVY,
              bar=BLUE, chip=BLUE):
    add_rect(s, x, y, w, h, fill=fill, rounded=True, radius=0.10, shadow=True)
    add_rect(s, x, y, 0.10, h, fill=bar, rounded=False, shadow=False)
    add_chip(s, x + 0.42, y + 0.33, 0.44, str(num), chip_color=chip,
             glyph_color=WHITE, gsize=13)
    blocks = [{"text": title, "size": 11, "bold": True, "color": title_color,
               "align": PP_ALIGN.LEFT, "space_after": 2}]
    for sline in subs:
        blocks.append({"text": sline, "size": 7.8, "color": GRAY,
                       "align": PP_ALIGN.LEFT, "space_after": 0,
                       "line_spacing": 1.0})
    tb = s.shapes.add_textbox(Inches(x + 0.70), Inches(y + 0.03),
                              Inches(w - 0.80), Inches(h - 0.06))
    _set_text(tb, blocks, anchor=MSO_ANCHOR.MIDDLE, l=0.02, r=0.04)

def stop_node(x, y, w, text):
    b = add_rect(s, x, y, w, 0.42, fill=STOP_SF, line=STOP, line_w=1.0,
                 rounded=True, radius=0.28, shadow=True)
    _set_text(b, [{"text": text, "size": 7.8, "bold": True, "color": STOP,
                   "align": PP_ALIGN.CENTER}])

def yn(x, y, text, color):
    add_text(s, x, y, 0.55, 0.24, [{"text": text, "size": 8.5, "bold": True,
             "color": color, "align": PP_ALIGN.CENTER}],
             anchor=MSO_ANCHOR.MIDDLE, l=0.01, r=0.01, t=0, b=0)

TOP = 1.5
LX, LW = 0.45, 3.05
cx1 = LX + LW / 2
stage_box(LX, TOP+0.00, LW, 0.80, 1, "Requisição de Vaga",
          ["Gestor solicita a abertura da vaga",
           "Cargo · Centro de custo · Escopo"])
stage_box(LX, TOP+0.94, LW, 0.80, 2, "Publicação da Vaga",
          ["Publicação após requisição aprovada",
           "Divulgação em painel e canais internos"])
stage_box(LX, TOP+1.88, LW, 0.74, 3, "Triagem de Currículos",
          ["Avaliação técnica dos currículos",
           "Seleção de candidatos aderentes"])
stage_box(LX, TOP+2.76, LW, 0.80, 4, "Entrevista RH",
          ["Avaliação comportamental · histórico",
           "Motivação · fit cultural · pretensão"])
for ya in (0.80, 1.74, 2.62):
    add_arrow(s, cx1, TOP+ya, cx1, TOP+ya+0.14, color=BLUE_MED, width=1.8)

d1cy = TOP + 4.35
add_arrow(s, cx1, TOP+3.56, cx1, d1cy-0.42, color=BLUE_MED, width=1.8)
add_diamond(s, cx1, d1cy, 1.5, 0.82, text="Aprovado?")
yn(LX-0.02, d1cy-0.34, "Não", STOP)
add_arrow(s, cx1-0.75, d1cy, LX+0.25, d1cy, color=STOP, width=1.5)
stop_node(LX-0.30, d1cy+0.30, 1.5, "Reprovado · feedback ao candidato")
yn(cx1+0.42, d1cy-0.02, "Sim", GREEN)

# Branch paralelo
PX, PW = 3.95, 2.05
pcx = PX + PW / 2
hdr = add_rect(s, PX, TOP, PW, 0.44, fill=YELLOW, rounded=True, radius=0.22,
               shadow=True)
_set_text(hdr, [{"text": "Triagem Patrimonial", "size": 10, "bold": True,
                 "color": NAVY, "align": PP_ALIGN.CENTER}])
add_text(s, PX, TOP+0.45, PW, 0.22, [{"text": "fluxo paralelo · background check",
         "size": 7.5, "italic": True, "color": GRAY, "align": PP_ALIGN.CENTER}],
         anchor=MSO_ANCHOR.TOP)
def mini(y, text):
    b = add_rect(s, PX, y, PW, 0.44, fill=WHITE, line=GRAY_LT, line_w=0.75,
                 rounded=True, radius=0.16, shadow=True)
    _set_text(b, [{"text": text, "size": 8.3, "bold": True, "color": GRAY_DK,
                   "align": PP_ALIGN.CENTER}])
mini(TOP+0.78, "Envio de formulário(s)")
mini(TOP+1.34, "Preenchimento pelo candidato")
mini(TOP+1.90, "Pesquisa patrimonial")
for ya in (0.72, 1.28, 1.84):
    add_arrow(s, pcx, TOP+ya, pcx, TOP+ya+0.06, color=YELLOW_DK, width=1.5)
pdcy = TOP + 3.15
add_arrow(s, pcx, TOP+2.34, pcx, pdcy-0.40, color=YELLOW_DK, width=1.5)
add_diamond(s, pcx, pdcy, 1.5, 0.80, fill=BLUE_SOFT, line=BLUE_MED,
            text="Aprovado?")
yn(pcx-0.5, pdcy+0.40, "Não", STOP)
add_arrow(s, pcx, pdcy+0.40, pcx, pdcy+0.58, color=STOP, width=1.4)
stop_node(PX, pdcy+0.58, PW, "Restrição · processo encerrado")
yn(pcx+0.58, pdcy-0.30, "Sim", GREEN)
# ligação D1 Sim -> topo branch
add_arrow(s, cx1+0.75, d1cy, 3.72, d1cy, color=GREEN, width=1.5, head=False)
add_arrow(s, 3.72, d1cy, 3.72, TOP+0.22, color=GREEN, width=1.5, head=False)
add_arrow(s, 3.72, TOP+0.22, PX, TOP+0.22, color=GREEN, width=1.5, head=True)

# Lane 2
QX, QW = 6.35, 3.05
qcx = QX + QW / 2
stage_box(QX, TOP+0.00, QW, 0.80, 5, "Entrevista com Gestor",
          ["Conhecimento técnico · experiência",
           "Competências · fit com a equipe"])
CH2 = 6.12
add_arrow(s, pcx+0.75, pdcy, CH2, pdcy, color=GREEN, width=1.5, head=False)
add_arrow(s, CH2, pdcy, CH2, TOP+0.40, color=GREEN, width=1.5, head=False)
add_arrow(s, CH2, TOP+0.40, QX, TOP+0.40, color=GREEN, width=1.5)
d2cy = TOP + 1.55
add_arrow(s, qcx, TOP+0.80, qcx, d2cy-0.40, color=BLUE_MED, width=1.8)
add_diamond(s, qcx, d2cy, 1.7, 0.80, text="Candidato\nFinalista?")
yn(qcx-1.02, d2cy-0.02, "Não", STOP)
add_arrow(s, qcx-0.85, d2cy, 6.28, d2cy, color=STOP, width=1.4, head=False)
add_arrow(s, 6.28, d2cy, 6.28, TOP+5.35, color=STOP, width=1.4, head=False)
add_arrow(s, 6.28, TOP+5.35, 6.62, TOP+5.35, color=STOP, width=1.4)
b = add_rect(s, 6.62, TOP+5.12, 2.78, 0.5, fill=STOP_SF, line=STOP, line_w=1.0,
             rounded=True, radius=0.16, shadow=True)
_set_text(b, [{"text": "Não finalista → nova triagem / nova publicação · "
               "feedback ao candidato", "size": 7.6, "bold": True,
               "color": STOP, "align": PP_ALIGN.CENTER}])
yn(qcx+0.10, d2cy+0.40, "Sim", GREEN)
stage_box(QX, TOP+2.30, QW, 0.74, 6, "Carta Proposta",
          ["Aprovação RH Custos · Aprovação Gestor",
           "Formalização da oferta ao candidato"],
          bar=YELLOW_DK, chip=YELLOW_DK)
add_arrow(s, qcx, d2cy+0.40, qcx, TOP+2.30, color=BLUE_MED, width=1.8)
stage_box(QX, TOP+3.18, QW, 0.80, 7, "Formalização & Assinatura",
          ["Envio da proposta por e-mail",
           "Assinatura eletrônica · aceite"])
add_arrow(s, qcx, TOP+3.04, qcx, TOP+3.18, color=BLUE_MED, width=1.8)
stage_box(QX, TOP+4.12, QW, 0.80, 8, "Seleção Concluída",
          ["Onboarding do novo colaborador",
           "Integração à equipe e à empresa"],
          bar=GREEN, chip=GREEN, title_color="1B5E20", fill=GREEN_SF)
add_arrow(s, qcx, TOP+3.98, qcx, TOP+4.12, color=GREEN, width=1.8)

# Painéis à direita
def panel(x, y, w, h, title, items, title_bg):
    add_rect(s, x, y, w, h, fill=WHITE, rounded=True, radius=0.06, shadow=True)
    add_rect(s, x, y, w, 0.38, fill=title_bg, rounded=True, radius=0.10,
             shadow=False)
    add_rect(s, x, y+0.20, w, 0.18, fill=title_bg, rounded=False, shadow=False)
    add_rect(s, x, y, 0.08, 0.38, fill=YELLOW, rounded=False, shadow=False)
    _th = s.shapes.add_textbox(Inches(x+0.05), Inches(y), Inches(w-0.1),
                               Inches(0.38))
    _set_text(_th, [{"text": title, "size": 9.5, "bold": True, "color": WHITE,
                     "align": PP_ALIGN.CENTER}])
    blocks = [{"text": "•  " + it, "size": 8, "color": GRAY_DK,
               "align": PP_ALIGN.LEFT, "space_after": 1.8, "line_spacing": 1.0}
              for it in items]
    tb = s.shapes.add_textbox(Inches(x+0.14), Inches(y+0.44), Inches(w-0.24),
                              Inches(h-0.5))
    _set_text(tb, blocks, anchor=MSO_ANCHOR.TOP, l=0.02, r=0.02)

PNX, PNW = 9.72, 3.2
panel(PNX, TOP+0.00, PNW, 1.72, "SISTEMAS (comuns / variáveis)",
      ["Pandapé (gestão de vagas)", "Outlook · Teams",
       "Assinatura eletrônica (GOV.BR / DocuSign)",
       "Formulários (Forms) · canais internos"], BLUE)
panel(PNX, TOP+1.86, PNW, 1.66, "PONTOS DE DECISÃO",
      ["Aprovação da Requisição", "Aprovação RH Custos",
       "Aprovação Entrevista RH", "Aprovação Triagem Patrimonial",
       "Aprovação Gestor", "Aceite da Proposta"], BLUE)
panel(PNX, TOP+3.66, PNW, 1.55, "POSSÍVEIS ENCERRAMENTOS",
      ["Reprovação na Entrevista RH", "Restrição na Triagem Patrimonial",
       "Reprovação pelo Gestor", "Recusa da Proposta", "Reabertura da vaga"],
      STOP)
finalize(s)


# ==========================================================================
# SLIDE 5 — VISÃO EXECUTIVA (LINHA DO TEMPO)
# ==========================================================================
s = new_slide(WHITE)
header(s, "Processo Padrão — Visão Executiva", tag="Linha do tempo · 8 etapas")

stages = [
    ("1", "Requisição", "Abertura da vaga pelo gestor"),
    ("2", "Divulgação", "Publicação da vaga (Pandapé)"),
    ("3", "Triagem", "Análise de currículos"),
    ("4", "Entrevistas", "RH e gestor + triagem patrimonial"),
    ("5", "Aprovação", "Decisões e validações"),
    ("6", "Proposta", "Carta e negociação"),
    ("7", "Admissão", "Assinatura eletrônica"),
    ("8", "Onboarding", "Integração do colaborador"),
]
line_y = 3.7
x_start, x_end = 0.95, 12.4
seg = (x_end - x_start) / (len(stages) - 1)
add_rect(s, x_start, line_y-0.03, x_end-x_start, 0.06, fill=GRAY_LT,
         rounded=False, shadow=False)
for i, (num, title, desc) in enumerate(stages):
    cx = x_start + seg * i
    up = (i % 2 == 0)
    node_col = BLUE if up else YELLOW_DK
    card_h = 1.28
    if up:
        card_y = line_y - 0.55 - card_h
        add_arrow(s, cx, line_y-0.32, cx, card_y+card_h, color=node_col,
                  width=1.4, head=False)
    else:
        card_y = line_y + 0.55
        add_arrow(s, cx, line_y+0.32, cx, card_y, color=node_col, width=1.4,
                  head=False)
    cw = 1.42
    add_rect(s, cx-cw/2, card_y, cw, card_h, fill=WHITE, line=GRAY_LT,
             line_w=0.75, rounded=True, radius=0.10, shadow=True)
    add_rect(s, cx-cw/2, card_y, cw, 0.09, fill=node_col, rounded=False,
             shadow=False)
    tb = s.shapes.add_textbox(Inches(cx-cw/2+0.06), Inches(card_y+0.14),
                              Inches(cw-0.12), Inches(card_h-0.18))
    _set_text(tb, [{"text": title, "size": 12, "bold": True, "color": NAVY,
                    "align": PP_ALIGN.CENTER, "space_after": 3},
                   {"text": desc, "size": 8.2, "color": GRAY,
                    "align": PP_ALIGN.CENTER, "line_spacing": 1.0}],
              anchor=MSO_ANCHOR.TOP)
    add_chip(s, cx, line_y, 0.56, num, chip_color=node_col, glyph_color=WHITE,
             gsize=15)

add_text(s, 0.4, 5.7, 12.5, 0.4,
         [{"text": "Do pedido à integração — um fluxo estruturado, auditável e "
           "orientado a decisões, comum a todas as unidades.", "size": 12,
           "italic": True, "color": GRAY, "align": PP_ALIGN.CENTER}],
         anchor=MSO_ANCHOR.MIDDLE)
ind_y = 6.2
indicators = [("8", "etapas principais", BLUE),
              ("6", "pontos de decisão", YELLOW_DK),
              ("5", "possíveis encerramentos", STOP),
              ("5", "unidades consolidadas", GREEN)]
iw = 2.85
gap = (13.333 - iw*len(indicators)) / (len(indicators)+1)
for i, (val, lbl, col) in enumerate(indicators):
    x = gap + i*(iw+gap)
    add_rect(s, x, ind_y, iw, 0.78, fill=GRAY_BG, rounded=True, radius=0.12,
             shadow=True)
    add_rect(s, x, ind_y, 0.10, 0.78, fill=col, rounded=False, shadow=False)
    add_text(s, x+0.18, ind_y, 0.85, 0.78, [{"text": val, "size": 24,
             "bold": True, "color": col, "align": PP_ALIGN.CENTER}],
             anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, x+1.02, ind_y, iw-1.1, 0.78, [{"text": lbl, "size": 10.5,
             "bold": True, "color": GRAY_DK, "align": PP_ALIGN.LEFT}],
             anchor=MSO_ANCHOR.MIDDLE)
finalize(s)


# ==========================================================================
# SLIDE 6 — PROCESSO PADRÃO: ONBOARDING (10 ETAPAS)
# ==========================================================================
s = new_slide(GRAY_BG)
header(s, "Processo Padrão — Onboarding / Admissão", tag="PADRÃO · 10 etapas",
       subtitle="Sequência de referência (modelo Brasil). Sistemas de folha e "
       "programas de acompanhamento variam por unidade — ver particularidades.")

ob = [
    ("1", "Aprovação do candidato", "RH · Pandapé",
     "Movimentação para Contratação; envio do link de documentação."),
    ("2", "Agendamento do ASO", "RH · Saúde Ocupacional",
     "Encaminhamento dos dados; exame admissional agendado."),
    ("3", "Documentação", "RH",
     "Conferência de documentos; organização do Book do Colaborador."),
    ("4", "Checklist Pré-Admissional", "RH · conferências",
     "Documentação, ASO, sistema, uniforme, crachá, ponto, benefícios, NR-22…"),
    ("5", "Cadastro no sistema de folha", "RH · sistema local",
     "Pré-requisitos: documentação completa e ASO apto."),
    ("6", "Assinatura Contratual", "RH · DocuSign",
     "Contrato de Trabalho e Ficha de Registro."),
    ("7", "Comunicação Pré-Integração", "RH · sexta anterior",
     "Guia de Integração, local/horário, rota de ônibus, NR-22, acessos."),
    ("8", "Primeiro dia de Integração", "RH",
     "Uniforme; ponto, benefícios, seguro de vida, formulários obrigatórios."),
    ("9", "Padrinho & Jornada de Experiência", "RH · 45 dias",
     "Acompanhamento da adaptação, treinamentos e feedbacks."),
    ("10", "Encerramento do Onboarding", "RH · Gestor",
     "Reunião de acompanhamento; registro do encerramento da jornada."),
]
cols, rows_n = 5, 2
cw, ch = 2.42, 2.28
gx, gy = 0.18, 0.22
x0 = 0.4
y0 = 1.55
for i, (num, title, sistema, desc) in enumerate(ob):
    r = i // cols
    c = i % cols
    x = x0 + c * (cw + gx)
    y = y0 + r * (ch + gy)
    accent = BLUE if r == 0 else YELLOW_DK
    add_rect(s, x, y, cw, ch, fill=WHITE, rounded=True, radius=0.07,
             shadow=True)
    add_rect(s, x, y, cw, 0.5, fill=accent, rounded=True, radius=0.12,
             shadow=False)
    add_rect(s, x, y+0.28, cw, 0.22, fill=accent, rounded=False, shadow=False)
    add_chip(s, x+0.36, y+0.25, 0.42, num, chip_color=WHITE, glyph_color=accent,
             gsize=13)
    add_text(s, x+0.62, y, cw-0.68, 0.5, [{"text": title, "size": 9.5,
             "bold": True, "color": WHITE, "align": PP_ALIGN.LEFT}],
             anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, x+0.12, y+0.58, cw-0.24, 0.3, [{"text": sistema, "size": 8,
             "bold": True, "color": accent, "align": PP_ALIGN.LEFT}],
             anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, x+0.12, y+0.9, cw-0.24, ch-1.0, [{"text": desc, "size": 8.3,
             "color": GRAY_DK, "align": PP_ALIGN.LEFT, "line_spacing": 1.05}],
             anchor=MSO_ANCHOR.TOP)
finalize(s)


# ==========================================================================
# SLIDE 7 — PARTICULARIDADES POR UNIDADE: R&S (MATRIZ)
# ==========================================================================
s = new_slide(WHITE)
header(s, "Particularidades por Unidade — Recrutamento e Seleção",
       tag="PARTICULARIDADES · rastreabilidade",
       subtitle="Cada célula documenta a exceção local; o cabeçalho identifica "
       "a unidade de origem. O que não diverge segue o processo padrão.")

rs_rows = [
    ["Dimensão do processo", "Aura Apoena", "Aura Borborema (BBR)",
     "Aura Serra Grande"],
    ["Origem da requisição",
     "Pandapé (Almas) / Portal de Serviços–HUB RH (Fluxo Original)",
     "Pandapé — aprovação: Gestor da área + Coord. de RH",
     "DocuSign"],
    ["Publicação / divulgação",
     "Pandapé + Painel de Vagas, Pixel Midea, Aura Comunica",
     "Pandapé; vagas estratégicas de R&S interno na CI Aura",
     "Pandapé sob demanda (7 dias); cartaz por e-mail e WhatsApp Business"],
    ["Triagem patrimonial",
     "Triagem Patrimonial + Formulário de Conflito de Interesses (questionário)",
     "“Background Check”: Patrimonial + Conflito de Interesses (2 formulários no Pandapé)",
     "Solicitação à Patrimonial por e-mail; ex-colaborador: histórico ao DP"],
    ["Entrevista com gestor",
     "Cargos SR+: superior imediato e gestor da gerência",
     "Cargos SR+: superior imediato e gestor da gerência",
     "Cargos SR+: superior imediato e gerentes"],
    ["Avaliação comportamental",
     "Não prevista explicitamente",
     "Não prevista explicitamente",
     "Etapa complementar para cargos sênior/gestão/críticos"],
    ["Assinatura da proposta",
     "GOV.BR (Fluxo Original) / DocuSign (Almas)",
     "DocuSign; RH Custos por e-mail, Gestor no Pandapé",
     "GOV.BR ou DocuSign (carta assinada)"],
    ["Sistemas específicos",
     "Teams, Outlook, Pandapé, Forms, DocuSign, GOV.BR, Pixel Midea, Aura Comunica",
     "Pandapé, DocuSign, Outlook (conjunto reduzido)",
     "Gen.te Mobile, Teams, Outlook, Pandapé, DocuSign, GOV.BR, WhatsApp Business"],
]
cell_colors = {}
for r in range(1, len(rs_rows)):
    cell_colors[(r, 1)] = BLUE_PALE
    cell_colors[(r, 2)] = TEAL_SF
    cell_colors[(r, 3)] = YELLOW_SF
add_table(s, 0.4, 1.5, 12.53, [2.35, 3.4, 3.4, 3.38], rs_rows,
          row_h=0.66, header_size=9.5, font_size=8.0, zebra=False,
          cell_colors=cell_colors)
# selos de rastreabilidade sob o cabeçalho não são necessários — colunas já nomeiam
add_text(s, 0.4, 6.95, 12.5, 0.3,
         [{"text": "Nota: o modelo de Honduras e Guatemala (LATAM) segue lógica "
           "distinta e é detalhado nos slides seguintes.", "size": 8.5,
           "italic": True, "color": GRAY, "align": PP_ALIGN.LEFT}],
         anchor=MSO_ANCHOR.MIDDLE)
finalize(s)


# ==========================================================================
# SLIDE 8 — PARTICULARIDADES POR UNIDADE: ONBOARDING (MATRIZ)
# ==========================================================================
s = new_slide(WHITE)
header(s, "Particularidades por Unidade — Onboarding / Admissão",
       tag="PARTICULARIDADES · rastreabilidade",
       subtitle="Diferenças em relação ao onboarding padrão de 10 etapas. "
       "Origem indicada em cada coluna.")

ob_rows = [
    ["Dimensão", "Aura Apoena", "Aura Borborema (BBR)", "Aura Serra Grande"],
    ["Sistema de folha (etapa 5)",
     "Linx (executado pelo RH)",
     "Linx via Vexia (terceirizado); chamado aberto no Pipefy (etapa 4)",
     "FPW / Forponto (executado pelo Departamento Pessoal – DP)"],
    ["Execução da admissão",
     "RH executa internamente",
     "Vexia executa; RH abre chamado e recebe documentação/benefícios",
     "DP executa o cadastro; documentação enviada ao DP"],
    ["Integração de mineração",
     "NR-22 na integração do 1º dia",
     "2º e 3º dia: Introdutório de Mineração (SSMA/Planta/Geotecnia/Patrimonial) – NR-22",
     "NR-22 na integração; conferência no checklist"],
    ["Programa de acompanhamento",
     "Programa Padrinho & Jornada de Experiência (45 dias)",
     "Capacitação na função (NRs/on the job) + Padrinho; avaliação em Forms (45 dias, gestor)",
     "“Programa On The Job – Área” (em revisão no momento)"],
    ["Comunicação pré-integração",
     "Guia de Boas-Vindas (mandala, rota, cronograma, horário) – quinta anterior",
     "Guia por aplicativo de mensagem – sexta anterior",
     "Guia por e-mail – sexta anterior"],
    ["Diferenciais registrados",
     "Kit Boas-vindas; ASO por Saúde Ocupacional",
     "Código de conduta no 1º dia; contrato digital",
     "ASO via planilha compartilhada; encerramento com formulário assinado por ambas as partes"],
]
cc = {}
for r in range(1, len(ob_rows)):
    cc[(r, 1)] = BLUE_PALE
    cc[(r, 2)] = TEAL_SF
    cc[(r, 3)] = YELLOW_SF
add_table(s, 0.4, 1.5, 12.53, [2.35, 3.4, 3.4, 3.38], ob_rows,
          row_h=0.74, header_size=9.5, font_size=8.0, zebra=False,
          cell_colors=cc)
finalize(s)


# ==========================================================================
# SLIDE 9 — MODELO LATAM: HONDURAS (AURA MINOSA) — 3 ROTAS
# ==========================================================================
s = new_slide(GRAY_BG)
header(s, "Particularidade Estrutural — Honduras (Aura Minosa)",
       tag="PARTICULARIDADE · modelo LATAM",
       subtitle="Modelo próprio: 3 rotas de recrutamento conforme a origem do "
       "candidato. Difere do fluxo único do Brasil e é mantido separado do padrão.")

routes = [
    ("8A3A3A", "RECRUTAMENTO INTERNO", "9 etapas",
     ["Solicitud de Personal (Fresh Service)",
      "Publicação da vaga por e-mail ao pessoal",
      "Decisão: há candidatos internos?",
      "Recepção de aplicações na RRHH",
      "Entrevistas com o líder da área",
      "Seleção com o encarregado (cuadro comparativo em Excel)",
      "Reclasificación por Fresh Service",
      "Indução ao posto via SST (por e-mail)"]),
    ("A85A2A", "EXTERNO EM COMUNIDADES", "13 etapas",
     ["Solicitud de Personal (Fresh Service)",
      "Compartilhar vaga com comunidades (e-mail)",
      "Recepção de perfis por comunidades",
      "Reclutamiento externo se não houver candidatos",
      "Entrevistas com o encarregado da área",
      "Avaliações psicométricas (PsicoSmart) e/ou técnicas",
      "Visto bueno de Relaciones Comunitarias",
      "Exames de saúde (clínica ocupacional) — decisão",
      "Indução e processo de contratação"]),
    ("6B4C9A", "RECRUTAMENTO EXTERNO", "15 etapas",
     ["Solicitud de Personal (Fresh Service)",
      "Busca de perfis no LinkedIn",
      "Seleção de currículos e contato com candidatos",
      "Entrevistas com RH + psicométricas/técnicas",
      "Detalhe da entrevista (Excel); resumo ao líder",
      "Entrevista com o líder da área",
      "Oferta laboral por DocuSign — decisão: aceita?",
      "Documentação e exames de saúde — decisão",
      "2º candidato mais apto se reprovado/recusado"]),
]
cw = 4.02
gx = 0.24
x0 = 0.4
y0 = 1.62
for i, (col, title, badge, steps) in enumerate(routes):
    x = x0 + i * (cw + gx)
    add_rect(s, x, y0, cw, 4.75, fill=WHITE, rounded=True, radius=0.05,
             shadow=True)
    add_rect(s, x, y0, cw, 0.72, fill=col, rounded=True, radius=0.08,
             shadow=False)
    add_rect(s, x, y0+0.4, cw, 0.32, fill=col, rounded=False, shadow=False)
    add_text(s, x+0.15, y0+0.05, cw-0.3, 0.44, [{"text": title, "size": 12.5,
             "bold": True, "color": WHITE, "align": PP_ALIGN.LEFT}],
             anchor=MSO_ANCHOR.MIDDLE)
    bdg = add_rect(s, x+0.15, y0+0.44, 1.3, 0.24, fill=WHITE, rounded=True,
                   radius=0.5, shadow=False)
    _set_text(bdg, [{"text": badge, "size": 8, "bold": True, "color": col,
                     "align": PP_ALIGN.CENTER}], l=0.02, r=0.02, t=0, b=0)
    blocks = []
    for j, st in enumerate(steps, 1):
        blocks.append({"runs": [
            {"text": f"{j}. ", "bold": True, "color": col, "size": 8.6},
            {"text": st, "color": GRAY_DK, "size": 8.6}],
            "align": PP_ALIGN.LEFT, "space_after": 3.2, "line_spacing": 1.02})
    add_text(s, x+0.16, y0+0.86, cw-0.3, 3.8, blocks, anchor=MSO_ANCHOR.TOP)

add_rect(s, 0.4, 6.5, 12.53, 0.62, fill=NAVY, rounded=True, radius=0.1,
         shadow=True)
add_text(s, 0.6, 6.5, 12.1, 0.62,
         [{"runs": [
             {"text": "Sistemas próprios: ", "bold": True, "color": YELLOW,
              "size": 9.5},
             {"text": "Fresh Service · PsicoSmart · LinkedIn · DocuSign · "
              "Excel · Correo. ", "color": WHITE, "size": 9.5},
             {"text": "Diferenciais: ", "bold": True, "color": YELLOW,
              "size": 9.5},
             {"text": "avaliações psicométricas, visto bueno de Relaciones "
              "Comunitarias e exames de saúde como ponto de decisão.",
              "color": WHITE, "size": 9.5}],
           "align": PP_ALIGN.LEFT}], anchor=MSO_ANCHOR.MIDDLE)
finalize(s)


# ==========================================================================
# SLIDE 10 — MODELO LATAM: HONDURAS ONBOARDING + GUATEMALA
# ==========================================================================
s = new_slide(WHITE)
header(s, "Particularidade Estrutural — Onboarding LATAM",
       tag="PARTICULARIDADE · Honduras & Guatemala",
       subtitle="Incorporación (Honduras) e Inducción (Guatemala) preservam "
       "etapas próprias não presentes no modelo Brasil.")

# Honduras onboarding (esquerda)
add_rect(s, 0.4, 1.55, 6.15, 5.25, fill=GRAY_BG, rounded=True, radius=0.04,
         shadow=True)
add_rect(s, 0.4, 1.55, 6.15, 0.5, fill="8A3A3A", rounded=True, radius=0.08,
         shadow=False)
add_rect(s, 0.4, 1.83, 6.15, 0.22, fill="8A3A3A", rounded=False, shadow=False)
add_text(s, 0.6, 1.55, 5.8, 0.5, [{"text": "Honduras (Minosa) — Incorporación · "
         "10 etapas", "size": 12, "bold": True, "color": WHITE,
         "align": PP_ALIGN.LEFT}], anchor=MSO_ANCHOR.MIDDLE)
hond = [
    "1. Confirmación de ingreso (RH)",
    "2. Inducción General – 2 dias (RH, SSMA, Relaciones Comunitarias)",
    "3. Firma e recepção de documentação → expediente físico",
    "4. Creación del expediente digital",
    "5. Asignación de recursos: uniforme, computadora, acessórios",
    "6. Registro en plataformas institucionales",
    "7. Inscripción en Clínica Ocupacional + seguros",
    "8. Comunicación al líder inmediato",
    "9. Inducción al puesto (líder; formato de inducción)",
    "10. Seguimiento del período de prueba (formato de evaluación)",
]
blocks = [{"text": h, "size": 9.2, "color": GRAY_DK, "align": PP_ALIGN.LEFT,
           "space_after": 4.2, "line_spacing": 1.03} for h in hond]
add_text(s, 0.62, 2.2, 5.7, 4.5, blocks, anchor=MSO_ANCHOR.TOP)

# Guatemala onboarding (direita)
add_rect(s, 6.78, 1.55, 6.15, 5.25, fill=GRAY_BG, rounded=True, radius=0.04,
         shadow=True)
add_rect(s, 6.78, 1.55, 6.15, 0.5, fill=PURPLE, rounded=True, radius=0.08,
         shadow=False)
add_rect(s, 6.78, 1.83, 6.15, 0.22, fill=PURPLE, rounded=False, shadow=False)
add_text(s, 6.98, 1.55, 5.8, 0.5, [{"text": "Guatemala — Inducción de Personal",
         "size": 12, "bold": True, "color": WHITE, "align": PP_ALIGN.LEFT}],
         anchor=MSO_ANCHOR.MIDDLE)
guat = [
    "Programação da Indução",
    "Requisição de ferramentas, equipamento, transporte e alimentação",
    "   Áreas: SSMA, Ambiente, Social, Seg. Patrimonial, Nómina",
    "Documentação: contrato, EPP, acordos, entrega de equipamento",
    "Inducción General e Inducción al Puesto (Forms, Outlook, Singular)",
    "Evaluación do Programa de Inducción",
    "Plan de Acción para cierre de brechas (quando necessário)",
    "Confirmación al puesto de trabajo",
]
blocks = [{"text": ("•  " + g if not g.startswith("   ") else g), "size": 9.2,
           "color": GRAY_DK, "align": PP_ALIGN.LEFT, "space_after": 5.5,
           "line_spacing": 1.03} for g in guat]
add_text(s, 7.0, 2.2, 5.7, 3.4, blocks, anchor=MSO_ANCHOR.TOP)
# aviso de lacuna Guatemala
add_rect(s, 7.0, 5.95, 5.7, 0.72, fill=PURPLE_SF, line=PURPLE, line_w=1.0,
         rounded=True, radius=0.12, shadow=False)
add_text(s, 7.15, 5.95, 5.4, 0.72, [{"runs": [
    {"text": "LACUNA: ", "bold": True, "color": PURPLE, "size": 9},
    {"text": "o fluxo de R&S (Procedimiento I e II) está apenas em imagens, "
     "sem texto extraível. Conteúdo não consolidável até ser fornecido.",
     "color": GRAY_DK, "size": 9}], "align": PP_ALIGN.LEFT,
    "line_spacing": 1.05}], anchor=MSO_ANCHOR.MIDDLE)
# sistemas Guatemala
add_text(s, 7.0, 5.6, 5.7, 0.3, [{"runs": [
    {"text": "Sistemas: ", "bold": True, "color": PURPLE, "size": 8.5},
    {"text": "FreshService, Outlook, Forms, Singular.", "color": GRAY,
     "size": 8.5}], "align": PP_ALIGN.LEFT}], anchor=MSO_ANCHOR.MIDDLE)
finalize(s)


# ==========================================================================
# SLIDE 11 — DIVERGÊNCIAS, INCONSISTÊNCIAS E CONFLITOS
# ==========================================================================
s = new_slide(WHITE)
header(s, "Divergências, Inconsistências e Conflitos",
       tag="Pontos que exigem decisão",
       subtitle="Onde as apresentações se contradizem ou descrevem a mesma "
       "etapa de formas incompatíveis.")

conf_rows = [
    ["#", "Ponto do processo", "Divergência / conflito observado", "Unidades"],
    ["1", "Origem da requisição",
     "Pandapé × Portal de Serviços (HUB RH) × DocuSign — inclusive divergência "
     "interna na própria Apoena (Almas usa Pandapé; Fluxo Original usa Portal).",
     "Apoena · BBR · Serra Grande"],
    ["2", "Assinatura da proposta/contrato",
     "GOV.BR × DocuSign × ambos — sem regra única de qual assinador é o oficial.",
     "Apoena · BBR · Serra Grande"],
    ["3", "Triagem patrimonial vs. Conflito de Interesses",
     "Nomenclatura “Triagem Patrimonial” × “Background Check”; Conflito de "
     "Interesses presente em Almas/BBR e ausente no Fluxo Original/base.",
     "Apoena · BBR · Serra Grande"],
    ["4", "Sistema de folha / admissão",
     "Linx (RH) × Linx via Vexia + Pipefy (terceirizado) × FPW/Forponto (DP) — "
     "sistemas e responsáveis diferentes para a mesma etapa.",
     "Apoena · BBR · Serra Grande"],
    ["5", "Avaliação comportamental",
     "Etapa formal apenas em Serra Grande (cargos sênior/críticos); ausente nas demais.",
     "Serra Grande"],
    ["6", "Modelo de recrutamento",
     "Fluxo único de 8 etapas (Brasil) × 3 rotas distintas por origem (Honduras) — "
     "estruturas de processo não equivalentes.",
     "Brasil × Honduras"],
    ["7", "Avaliação psicométrica",
     "Etapa obrigatória no modelo LATAM (PsicoSmart) e inexistente no Brasil.",
     "Honduras × Brasil"],
]
cc = {(r, 3): STOP_SF for r in range(1, len(conf_rows))}
add_table(s, 0.4, 1.5, 12.53, [0.5, 2.5, 6.9, 2.63], conf_rows,
          row_h=0.72, header_fill=STOP, header_size=9.5, font_size=8.2,
          zebra=True, first_col_bold=True, cell_colors=cc)
finalize(s)


# ==========================================================================
# SLIDE 12 — LACUNAS DE INFORMAÇÃO
# ==========================================================================
s = new_slide(GRAY_BG)
header(s, "Lacunas de Informação Identificadas",
       tag="Pendências de levantamento",
       subtitle="Pontos não documentados, incompletos ou “em revisão” que "
       "impedem o fechamento do padrão institucional.")

gaps = [
    ("Guatemala — R&S sem texto",
     "Procedimiento I e II estão apenas em imagens. O fluxo de R&S não pôde ser "
     "consolidado; requer versão textual/editável."),
    ("Serra Grande — programa em revisão",
     "O “Programa On The Job – Área” consta como “em revisão no momento”; a "
     "etapa 9 do onboarding ainda não está estabilizada."),
    ("Prazos / SLA não padronizados",
     "A maioria das unidades não documenta prazos por etapa (ex.: dias de "
     "publicação, tempo de triagem, prazo de assinatura)."),
    ("Pontos de decisão do modelo LATAM",
     "Honduras/Guatemala não explicitam aprovações formais (RH Custos, Gestor) "
     "nos moldes do Brasil; papéis de aprovação pouco descritos."),
    ("Papel do Pandapé não uniforme",
     "Pandapé aparece ora como origem da requisição, ora apenas como canal de "
     "publicação — função do sistema não é consistente entre unidades."),
    ("Conflito de Interesses incompleto",
     "Nem todas as unidades documentam o formulário de Conflito de Interesses; "
     "onde existe, o critério de reprovação não é detalhado."),
]
cols = 2
cw = 6.15
ch = 1.6
gx = 0.24
gy = 0.2
x0, y0 = 0.4, 1.55
for i, (title, desc) in enumerate(gaps):
    r = i // cols
    c = i % cols
    x = x0 + c * (cw + gx)
    y = y0 + r * (ch + gy)
    add_rect(s, x, y, cw, ch, fill=WHITE, rounded=True, radius=0.06,
             shadow=True)
    add_rect(s, x, y, 0.14, ch, fill=PURPLE, rounded=False, shadow=False)
    add_chip(s, x+0.5, y+0.4, 0.5, "!", chip_color=PURPLE, glyph_color=WHITE,
             gsize=15)
    add_text(s, x+0.85, y+0.12, cw-1.0, 0.5, [{"text": title, "size": 11.5,
             "bold": True, "color": NAVY, "align": PP_ALIGN.LEFT}],
             anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, x+0.85, y+0.6, cw-1.0, ch-0.72, [{"text": desc, "size": 9.5,
             "color": GRAY_DK, "align": PP_ALIGN.LEFT, "line_spacing": 1.06}],
             anchor=MSO_ANCHOR.TOP)
finalize(s)


# ==========================================================================
# SLIDE 13 — RESUMO DAS PRINCIPAIS DIFERENÇAS ENTRE UNIDADES
# ==========================================================================
s = new_slide(WHITE)
header(s, "Resumo — Principais Diferenças entre Unidades",
       tag="Visão comparativa",
       subtitle="Leitura rápida de onde cada unidade se afasta do padrão "
       "institucional.")

sum_rows = [
    ["Unidade / País", "Modelo", "Sistema de folha", "Assinatura",
     "Marca distintiva"],
    ["Aura Apoena · BR", "Fluxo padrão 8 etapas", "Linx (RH)",
     "GOV.BR / DocuSign", "Conflito de Interesses; Kit e Guia de Boas-vindas"],
    ["Aura Borborema · BR", "Fluxo padrão 8 etapas", "Linx via Vexia + Pipefy",
     "DocuSign", "Admissão terceirizada; Introdutório de Mineração (2º/3º dia)"],
    ["Aura Serra Grande · BR", "Fluxo padrão + av. comportamental",
     "FPW/Forponto (DP)", "GOV.BR / DocuSign",
     "Avaliação comportamental; DP no cadastro e histórico funcional"],
    ["Aura Minosa · Honduras", "3 rotas por origem", "Plataforma institucional",
     "DocuSign", "Psicométricas (PsicoSmart); visto bueno Relaciones Comunitarias"],
    ["Aura Guatemala · Guatemala", "R&S em imagens (lacuna)", "Não informado",
     "Não informado", "Indução com cierre de brechas; sistema Singular"],
]
cell_colors = {}
units_order = ["Apoena", "Borborema", "Serra Grande", "Honduras", "Guatemala"]
for r in range(1, len(sum_rows)):
    cell_colors[(r, 0)] = UNIT_COLORS[units_order[r-1]]
add_table(s, 0.4, 1.5, 12.53, [2.55, 2.75, 2.35, 1.85, 3.03], sum_rows,
          row_h=0.82, header_size=9.5, font_size=8.4, zebra=True,
          first_col_bold=True, cell_colors=cell_colors,
          white_text_cells={(r, 0) for r in range(1, len(sum_rows))})
# faixa de leitura
add_text(s, 0.4, 6.95, 12.5, 0.3, [{"text": "Os três negócios do Brasil "
         "compartilham o mesmo fluxo; as diferenças concentram-se em sistemas "
         "de folha, assinatura e programas de acompanhamento. A LATAM opera "
         "modelo próprio.", "size": 8.8, "italic": True, "color": GRAY,
         "align": PP_ALIGN.LEFT}], anchor=MSO_ANCHOR.MIDDLE)
finalize(s)


# ==========================================================================
# SLIDE 14 — OPORTUNIDADES DE PADRONIZAÇÃO
# ==========================================================================
s = new_slide(GRAY_BG)
header(s, "Oportunidades de Padronização",
       tag="Recomendações",
       subtitle="Onde a organização pode convergir para um padrão único, "
       "preservando exceções legítimas (ex.: exigências legais locais).")

opps = [
    ("1", "Origem única da requisição",
     "Definir um sistema oficial de abertura de vaga (ex.: Pandapé) e o fluxo "
     "de aprovação padrão (Gestor + RH), eliminando Portal/DocuSign como origem."),
    ("2", "Assinatura eletrônica unificada",
     "Eleger um assinador oficial (GOV.BR ou DocuSign) e usar o outro apenas "
     "como contingência, com regra clara."),
    ("3", "Background Check padronizado",
     "Unificar nomenclatura e escopo: Triagem Patrimonial + Conflito de "
     "Interesses obrigatórios em todas as unidades, com critérios de reprovação."),
    ("4", "Modelo de folha e responsáveis",
     "Padronizar sistema de folha e definir quem executa a admissão (RH, DP ou "
     "parceiro), com SLA e checklist único."),
    ("5", "Onboarding com base comum + módulos locais",
     "Manter as 10 etapas como núcleo e tratar mineração/comunidades como "
     "módulos locais oficiais, não como variações informais."),
    ("6", "Governança de prazos e documentação",
     "Definir SLA por etapa, completar o R&S de Guatemala e estabilizar o "
     "programa de acompanhamento de Serra Grande."),
]
cols = 2
cw = 6.15
ch = 1.6
gx, gy = 0.24, 0.2
x0, y0 = 0.4, 1.55
for i, (num, title, desc) in enumerate(opps):
    r = i // cols
    c = i % cols
    x = x0 + c * (cw + gx)
    y = y0 + r * (ch + gy)
    add_rect(s, x, y, cw, ch, fill=WHITE, rounded=True, radius=0.06,
             shadow=True)
    add_rect(s, x, y, 0.14, ch, fill=GREEN, rounded=False, shadow=False)
    add_chip(s, x+0.5, y+0.4, 0.5, num, chip_color=GREEN, glyph_color=WHITE,
             gsize=15)
    add_text(s, x+0.85, y+0.12, cw-1.0, 0.5, [{"text": title, "size": 11.5,
             "bold": True, "color": NAVY, "align": PP_ALIGN.LEFT}],
             anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, x+0.85, y+0.6, cw-1.0, ch-0.72, [{"text": desc, "size": 9.3,
             "color": GRAY_DK, "align": PP_ALIGN.LEFT, "line_spacing": 1.05}],
             anchor=MSO_ANCHOR.TOP)
finalize(s)


# ==========================================================================
OUT = "Fluxo_Recrutamento_Selecao_Consolidado.pptx"
prs.save(OUT)
print("OK ->", OUT, "| slides:", slide_no[0])
